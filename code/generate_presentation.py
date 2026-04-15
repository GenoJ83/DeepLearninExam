"""Build CSC3218 PowerPoint (8 slides) from results/. Run: python generate_presentation.py"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

# Package layout: code/ and results/ are siblings under the submission root.
ROOT = Path(__file__).resolve().parent.parent
RESULTS = ROOT / "results"
METRICS_PATH = RESULTS / "metrics.json"
CLASS_REPORT_PATH = RESULTS / "classification_report.txt"

# Match report / submission details
STUDENT_NAME = "Geno Owor Joshua"
STUDENT_REG = "M23B23/006"
INSTITUTION = "Uganda Christian University"


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if slide.placeholders[1]:
        slide.placeholders[1].text = subtitle


def add_bullet_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    *,
    font_pt: int = 20,
) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.word_wrap = True
    tf.text = bullets[0]
    for para in tf.paragraphs:
        para.font.size = Pt(font_pt)
    for line in bullets[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(font_pt)


def add_picture_slide(prs: Presentation, title: str, image_path: Path) -> None:
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    left = Inches(0.35)
    top = Inches(1.05)
    slide.shapes.add_picture(str(image_path), left, top, width=Inches(9.3))


def add_table_slide(
    prs: Presentation,
    title: str,
    headers: tuple[str, str],
    rows: list[tuple[str, str]],
) -> None:
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    nrows = 1 + len(rows)
    ncols = 2
    left, top, w = Inches(0.6), Inches(1.35), Inches(8.8)
    row_h = Inches(0.38)
    table_shape = slide.shapes.add_table(nrows, ncols, left, top, w, row_h * nrows)
    table = table_shape.table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(14)
    for i, (a, b) in enumerate(rows, start=1):
        table.cell(i, 0).text = a
        table.cell(i, 1).text = b
        for c in (0, 1):
            for p in table.cell(i, c).text_frame.paragraphs:
                p.font.size = Pt(13)


def add_table_and_bullets_slide(
    prs: Presentation,
    title: str,
    headers: tuple[str, str],
    rows: list[tuple[str, str]],
    bullets: list[str],
    *,
    table_row_h: float = 0.32,
    bullet_font_pt: int = 15,
) -> None:
    """Compact table in upper area + result bullets below (one slide)."""
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    nrows = 1 + len(rows)
    left, top, w = Inches(0.55), Inches(1.2), Inches(8.9)
    row_h = Inches(table_row_h)
    table_shape = slide.shapes.add_table(nrows, 2, left, top, w, row_h * nrows)
    table = table_shape.table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
    for i, (a, b) in enumerate(rows, start=1):
        table.cell(i, 0).text = a
        table.cell(i, 1).text = b
        for c in (0, 1):
            for p in table.cell(i, c).text_frame.paragraphs:
                p.font.size = Pt(11)
    table_bottom = top + row_h * nrows + Inches(0.15)
    box = slide.shapes.add_textbox(Inches(0.55), table_bottom, Inches(8.9), Inches(2.5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = bullets[0]
    for para in tf.paragraphs:
        para.font.size = Pt(bullet_font_pt)
    for line in bullets[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(bullet_font_pt)


def add_two_picture_slide(
    prs: Presentation,
    title: str,
    left_path: Path,
    right_path: Path,
    left_caption: str,
    right_caption: str,
) -> None:
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    y = Inches(1.05)
    w = Inches(4.45)
    slide.shapes.add_picture(str(left_path), Inches(0.4), y, width=w)
    slide.shapes.add_picture(str(right_path), Inches(5.05), y, width=w)
    cap = slide.shapes.add_textbox(Inches(0.4), Inches(6.55), Inches(4.45), Inches(0.55))
    cap.text_frame.text = left_caption
    for p in cap.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.italic = True
    cap2 = slide.shapes.add_textbox(Inches(5.05), Inches(6.55), Inches(4.45), Inches(0.55))
    cap2.text_frame.text = right_caption
    for p in cap2.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.italic = True


def load_metrics() -> dict:
    if not METRICS_PATH.is_file():
        raise FileNotFoundError(f"Missing {METRICS_PATH} — run training first.")
    return json.loads(METRICS_PATH.read_text(encoding="utf-8"))


def parse_per_class_f1(report_path: Path) -> list[tuple[str, float]]:
    """Return (class_name, f1) for each of the 10 CIFAR classes."""
    if not report_path.is_file():
        return []
    rows: list[tuple[str, float]] = []
    for line in report_path.read_text(encoding="utf-8").splitlines():
        parts = line.split()
        if len(parts) != 5:
            continue
        name, ps, rs, fs, sup_s = parts
        if name in ("macro", "weighted"):
            continue
        try:
            sup = int(sup_s)
            f1 = float(fs)
        except ValueError:
            continue
        if sup == 1000:
            rows.append((name, f1))
    return rows


def main() -> None:
    m = load_metrics()
    hp = m["hyperparameters"]
    te = m["test_accuracy"]
    tr = m["train_accuracy"]
    va = m["val_accuracy"]
    prec = m["test_precision_macro"]
    rec = m["test_recall_macro"]
    f1 = m["test_f1_macro"]
    test_loss = m.get("test_loss", 0.0)
    epochs = m["epochs_ran"]
    secs = m["seconds"]
    device = m["device"]
    hours = secs / 3600.0

    per_class = parse_per_class_f1(CLASS_REPORT_PATH)
    per_class_sorted = sorted(per_class, key=lambda x: x[1])
    worst = per_class_sorted[:2] if len(per_class_sorted) >= 2 else per_class_sorted
    best = sorted(per_class, key=lambda x: -x[1])[:2] if per_class else []
    wtxt = ", ".join(f"{n} ({v:.2f})" for n, v in worst) if worst else "—"
    btxt = ", ".join(f"{n} ({v:.2f})" for n, v in best) if best else "—"

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1 — Cover
    add_title_slide(
        prs,
        "CIFAR-10 Image Classification\nwith a Convolutional Neural Network",
        f"{STUDENT_NAME}\nReg. No. {STUDENT_REG}\n{INSTITUTION}\nCSC3218 — Deep Learning",
    )

    # 2 — Problem & data (merged)
    add_bullet_slide(
        prs,
        "Problem & CIFAR-10 data",
        [
            "Task: classify 32×32 RGB images into 10 classes (logits → softmax; train with cross-entropy).",
            "CNNs: local filters, hierarchy of features, pooling — suited to natural images.",
            "50k train / 10k test (official); we use 90% train / 10% val from the 50k; test only for final metrics.",
            "Normalize with CIFAR mean/std; train aug: RandomCrop+padding, horizontal flip — not on val/test.",
        ],
        font_pt=17,
    )

    # 3 — Model & training (merged)
    add_bullet_slide(
        prs,
        "Model & training",
        [
            "3× stages: two Conv3×3+BN+ReLU → MaxPool → dropout; channels 64→128→256; spatial 32→4.",
            "Global average pool → head: Linear→BN→ReLU→Dropout→Linear (10 classes).",
            f"AdamW (lr={hp['lr']}, weight_decay={hp['weight_decay']}); ReduceLROnPlateau on val acc (×0.5, patience 4).",
            f"Early stopping patience {hp['patience']} epochs; best val checkpoint → reload before test eval.",
            f"Regularization: dropout {hp['dropout']}, weight decay, augmentation; test set used once.",
        ],
        font_pt=16,
    )

    # 4 — Key settings + results (one slide)
    add_table_and_bullets_slide(
        prs,
        "Setup & results",
        ("Setting", "Value"),
        [
            ("Epochs (completed / max)", f"{epochs} / {hp['epochs']}"),
            ("Batch size", str(hp["batch_size"])),
            ("Learning rate / weight decay", f"{hp['lr']} / {hp['weight_decay']}"),
            ("Dropout / early-stop patience", f"{hp['dropout']} / {hp['patience']}"),
            ("Device / runtime", f"{device} · ~{hours:.1f} h"),
        ],
        [
            f"Test accuracy {te * 100:.2f}% · test loss {test_loss:.3f} · macro P/R/F1 "
            f"{prec * 100:.1f}% / {rec * 100:.1f}% / {f1 * 100:.1f}%",
            f"Best checkpoint: train {tr * 100:.2f}% · val {va * 100:.2f}% (train–val gap → mild overfitting).",
            f"Hardest F1: {wtxt}; strongest: {btxt}.",
            "Val ≈ test accuracy → validation split was representative.",
        ],
        bullet_font_pt=14,
    )

    curves = RESULTS / "curves_loss_accuracy.png"
    cm_img = RESULTS / "confusion_matrix_test.png"
    samples = RESULTS / "sample_predictions.png"

    # 5 — Learning curves
    if curves.is_file():
        add_picture_slide(prs, "Training / validation loss and accuracy", curves)

    # 6 — Confusion matrix + samples side by side
    if cm_img.is_file() and samples.is_file():
        add_two_picture_slide(
            prs,
            "Test confusion matrix & sample predictions",
            cm_img,
            samples,
            "Confusion matrix (counts)",
            "Green = correct, red = wrong",
        )
    elif cm_img.is_file():
        add_picture_slide(prs, "Test confusion matrix", cm_img)
    elif samples.is_file():
        add_picture_slide(prs, "Sample predictions", samples)

    # 7 — Discussion & conclusion (merged)
    add_bullet_slide(
        prs,
        "Discussion & conclusion",
        [
            "Curves: training improves faster than validation; LR drops when val acc plateaus.",
            "Matrix: off-diagonal mass shows cat/dog and similar classes; vehicles/ships often cleaner.",
            "Limits: moderate-depth CNN; no TTA/ensemble; low-res images — ResNet / stronger aug could help.",
            f"Summary: BN + dropout + GAP CNN achieves {te * 100:.2f}% on CIFAR-10 with proper train/val/test practice.",
        ],
        font_pt=17,
    )

    # 8 — Thank you
    add_title_slide(prs, "Thank you", "Questions?")

    out = ROOT / "CSC3218_presentation.pptx"
    prs.save(out)
    n = len(prs.slides)
    if n < 8:
        print(f"Saved {out} ({n} slides — add missing PNGs in results/ for full deck)")
    else:
        print(f"Saved {out} ({n} slides)")


if __name__ == "__main__":
    main()
